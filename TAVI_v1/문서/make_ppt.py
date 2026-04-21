"""
학회 발표 PPT 생성
경희대학교 기후사회과학융합학과 / 도시컴퓨팅연구실
13분 발표용 (~15 슬라이드)
"""

import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import pptx.util as util

BASE    = os.path.dirname(os.path.abspath(__file__))
FIG_DIR = os.path.join(BASE, '../04_분석결과/figures')
OUT     = os.path.join(BASE, 'TAVI_presentation.pptx')

# ── 색상 팔레트 ──────────────────────────────────────────────────────
C_NAVY   = RGBColor(0x1A, 0x23, 0x7E)   # 진한 남색 (제목 배경)
C_BLUE   = RGBColor(0x1E, 0x88, 0xE5)   # 포인트 파란색
C_RED    = RGBColor(0xE5, 0x39, 0x35)   # 강조 빨간색
C_GREEN  = RGBColor(0x2E, 0x7D, 0x32)   # 녹색
C_GRAY   = RGBColor(0x42, 0x42, 0x42)   # 본문 회색
C_LGRAY  = RGBColor(0xF5, 0xF5, 0xF5)   # 배경 연회색
C_WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
C_ORANGE = RGBColor(0xEF, 0x6C, 0x00)
C_ACCENT = RGBColor(0x00, 0xAC, 0xC1)   # 청록

W  = Inches(13.33)   # 와이드 16:9
H  = Inches(7.5)

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H

BLANK = prs.slide_layouts[6]   # 완전 빈 슬라이드


# ── 헬퍼 함수 ─────────────────────────────────────────────────────────
def add_rect(slide, x, y, w, h, fill=None, line=None, alpha=None):
    shape = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.line.fill.background() if line is None else None
    if fill:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    else:
        shape.fill.background()
    if line is None:
        shape.line.fill.background()
    return shape

def add_text(slide, text, x, y, w, h, size=18, bold=False, color=C_GRAY,
             align=PP_ALIGN.LEFT, wrap=True, italic=False):
    txBox = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return txBox

def add_img(slide, path, x, y, w, h=None):
    if not os.path.exists(path):
        return
    if h:
        slide.shapes.add_picture(path, Inches(x), Inches(y), Inches(w), Inches(h))
    else:
        slide.shapes.add_picture(path, Inches(x), Inches(y), Inches(w))

def slide_header(slide, title, subtitle=None, bar_color=C_NAVY):
    """상단 컬러 바 + 제목"""
    add_rect(slide, 0, 0, 13.33, 1.2, fill=bar_color)
    add_text(slide, title, 0.4, 0.15, 12.0, 0.8,
             size=28, bold=True, color=C_WHITE, align=PP_ALIGN.LEFT)
    if subtitle:
        add_text(slide, subtitle, 0.4, 0.85, 12.0, 0.4,
                 size=13, color=RGBColor(0xBB, 0xDE, 0xFB), align=PP_ALIGN.LEFT)
    # 하단 라인
    add_rect(slide, 0, 7.3, 13.33, 0.2, fill=bar_color)
    # 슬라이드 번호 영역 (하단 우측)
    return slide

def bullet(slide, items, x, y, w, h_each=0.42, start_y=None,
           size=15, color=C_GRAY, marker="▪ ", bold_first=False):
    """불릿 리스트"""
    sy = start_y if start_y else y
    for i, item in enumerate(items):
        txt = marker + item if not item.startswith("  ") else item
        is_sub = item.startswith("  ")
        add_text(slide, txt, x + (0.3 if is_sub else 0),
                 sy + i * h_each, w - (0.3 if is_sub else 0),
                 h_each, size=size - (1 if is_sub else 0),
                 color=color, bold=(i == 0 and bold_first))


# ════════════════════════════════════════════════════════════════════
# Slide 1: 표지
# ════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)

# 배경
add_rect(sl, 0, 0, 13.33, 7.5, fill=C_NAVY)
add_rect(sl, 0, 4.8, 13.33, 2.7, fill=RGBColor(0x0D, 0x47, 0xA1))

# 포인트 라인
add_rect(sl, 0.5, 1.6, 0.08, 3.5, fill=C_BLUE)
add_rect(sl, 0.5, 1.6, 5.0,  0.06, fill=C_BLUE)

# 메인 타이틀
add_text(sl,
    "열환경을 반영한 보행 접근권역 분석",
    0.7, 1.7, 11.5, 1.1, size=36, bold=True, color=C_WHITE)
add_text(sl,
    "Thermal Catchment 고도화: UTCI 기반 링크 회피 모델",
    0.7, 2.7, 11.5, 0.7, size=22, color=RGBColor(0x90, 0xCA, 0xF9))
add_text(sl,
    "— 서울 성동구 7개 지하철역을 중심으로 —",
    0.7, 3.3, 11.5, 0.5, size=16, color=RGBColor(0xBB, 0xDE, 0xFB), italic=True)

# 발표자 정보
add_text(sl, "신 진", 1.0, 5.1, 5.0, 0.6, size=20, bold=True, color=C_WHITE)
add_text(sl, "경희대학교 기후사회과학융합학과  |  도시컴퓨팅연구실",
         1.0, 5.65, 9.0, 0.45, size=13, color=RGBColor(0xBB, 0xDE, 0xFB))
add_text(sl, "2026", 1.0, 6.1, 3.0, 0.4, size=13,
         color=RGBColor(0x90, 0xCA, 0xF9))


# ════════════════════════════════════════════════════════════════════
# Slide 2: 연구 배경
# ════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
add_rect(sl, 0, 0, 13.33, 7.5, fill=C_LGRAY)
slide_header(sl, "연구 배경", "왜 폭염과 보행 접근성인가?")

# 왼쪽 박스
add_rect(sl, 0.4, 1.4, 5.8, 4.5, fill=C_WHITE)
add_rect(sl, 0.4, 1.4, 0.1, 4.5, fill=C_RED)
add_text(sl, "🌡  폭염의 심각성", 0.65, 1.5, 5.3, 0.5, size=16, bold=True, color=C_RED)
bullet(sl, [
    "2025년 서울 폭염 일수 역대 최다",
    "UTCI 38°C 이상 = '매우 강한 열스트레스'",
    "  (Bröde et al. 2012, 국제 표준)",
    "야외 보행 중 열사병·탈수 위험 급증",
    "고령자·저소득층 집중 피해",
], 0.65, 2.05, 5.3, h_each=0.44, size=14, color=C_GRAY)

# 오른쪽 박스
add_rect(sl, 6.7, 1.4, 6.2, 4.5, fill=C_WHITE)
add_rect(sl, 6.7, 1.4, 0.1, 4.5, fill=C_BLUE)
add_text(sl, "🚶  보행 접근성의 공백", 6.95, 1.5, 5.7, 0.5, size=16, bold=True, color=C_BLUE)
bullet(sl, [
    "기존 접근성 연구: 거리·시간만 고려",
    "열환경(UTCI)을 소프트 제약으로 본",
    "  연구는 극히 드뭄",
    "폭염 = 일부 경로가 사실상 '통행 불가'",
    "→ 취약 집단의 이동권 불평등 심화",
], 6.95, 2.05, 5.7, h_each=0.44, size=14, color=C_GRAY)

# 하단 핵심 질문
add_rect(sl, 0.4, 6.1, 12.5, 1.05, fill=RGBColor(0xE3, 0xF2, 0xFD))
add_rect(sl, 0.4, 6.1, 0.12, 1.05, fill=C_BLUE)
add_text(sl,
    '연구 질문:  "폭염 시 더운 길을 회피하면, 지하철역까지의 15분 도달 가능 범위(Catchment)는 얼마나 줄어드는가?"',
    0.65, 6.2, 12.1, 0.8, size=15, bold=True, color=C_NAVY, align=PP_ALIGN.CENTER)


# ════════════════════════════════════════════════════════════════════
# Slide 3: 선행연구 및 연구 위치
# ════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
add_rect(sl, 0, 0, 13.33, 7.5, fill=C_LGRAY)
slide_header(sl, "선행연구 및 본 연구의 위치")

# 3개 카드
cards = [
    (C_NAVY,   "Catchment / 접근권역",
     ["15분 도시 (Moreno et al. 2021): 시간예산 기준",
      "Dijkstra 기반 도달 범위 = Catchment",
      "기존 연구: 거리·시간만 고려 (열환경 無)"]),
    (C_BLUE,   "UTCI & 열환경",
     ["Bröde et al. (2012): UTCI 국제 표준 카테고리",
      "Oke (1987): 도시 협곡 SVF 공식",
      "Lindberg & Grimmond (2011): SVF-MRT 관계"]),
    (C_GREEN,  "선행 유사 연구",
     ["Colaninno et al. (2024): SOLWEIG + LiDAR",
      "→ 보도 세그먼트 열위험 평가 (뉴욕)",
      "본 연구: 합성 DSM + Catchment 고도화"]),
]
for i, (col, title, items) in enumerate(cards):
    x0 = 0.4 + i * 4.3
    add_rect(sl, x0, 1.4, 4.1, 5.1, fill=C_WHITE)
    add_rect(sl, x0, 1.4, 4.1, 0.55, fill=col)
    add_text(sl, title, x0 + 0.15, 1.46, 3.8, 0.45,
             size=15, bold=True, color=C_WHITE)
    for j, item in enumerate(items):
        add_text(sl, ("▪ " if not item.startswith("→") else "") + item,
                 x0 + 0.2, 2.1 + j * 0.65, 3.8, 0.6,
                 size=13, color=C_GRAY)

# 본 연구 위치 강조
add_rect(sl, 0.4, 6.55, 12.5, 0.8, fill=RGBColor(0xE8, 0xF5, 0xE9))
add_rect(sl, 0.4, 6.55, 0.12, 0.8, fill=C_GREEN)
add_text(sl,
    "본 연구: 기존 Catchment에 열환경(UTCI) 제약을 추가한 Thermal Catchment 고도화 → 향후 STP로 확장 예정",
    0.65, 6.65, 12.1, 0.6, size=13, bold=True, color=C_NAVY, align=PP_ALIGN.CENTER)


# ════════════════════════════════════════════════════════════════════
# Slide 4: 연구 방법 개요
# ════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
add_rect(sl, 0, 0, 13.33, 7.5, fill=C_LGRAY)
slide_header(sl, "연구 방법 개요", "5단계 분석 파이프라인")

steps = [
    (C_NAVY,   "① 데이터 수집",
     "S-DoT 센서(57개) + ASOS\nOSM 보행 네트워크\n건물 폴리곤 (층수 속성)"),
    (C_BLUE,   "② UTCI 계산",
     "Open-Meteo 기상 취득\nSOLWEIG 방식 MRT 계산\npythermalcomfort → UTCI"),
    (C_ACCENT, "③ SVF 보정",
     "Oke(1987) H/W 공식\n링크 버퍼 20m 건물 높이\n캐노피 비율 계산"),
    (C_GREEN,  "④ 링크 회피",
     "UTCI ≥ 38°C 링크 제거\n(Bröde et al. 2012)\nHard-cut 모델"),
    (C_ORANGE, "⑤ Thermal Catchment",
     "Dijkstra 15분 도달 범위\nClassic vs Thermal 비교\n역별 감소율 산출"),
]

for i, (col, title, desc) in enumerate(steps):
    x0 = 0.35 + i * 2.55
    add_rect(sl, x0, 1.4, 2.3, 4.8, fill=C_WHITE)
    add_rect(sl, x0, 1.4, 2.3, 0.7, fill=col)
    add_text(sl, title, x0 + 0.12, 1.47, 2.1, 0.58,
             size=14, bold=True, color=C_WHITE)
    add_text(sl, desc, x0 + 0.12, 2.2, 2.1, 2.8,
             size=12.5, color=C_GRAY)
    # 화살표 (마지막 제외)
    if i < 4:
        add_text(sl, "→", x0 + 2.25, 2.7, 0.4, 0.5,
                 size=20, bold=True, color=C_NAVY, align=PP_ALIGN.CENTER)

# 하단 출력
add_rect(sl, 0.35, 6.35, 12.6, 0.9, fill=RGBColor(0xFF, 0xF8, 0xE1))
add_rect(sl, 0.35, 6.35, 0.12, 0.9, fill=C_ORANGE)
add_text(sl,
    "출력:  역별 접근성 감소율(%) + 결정 요인 회귀분석  |  연구지역: 성동구 7개 역  |  분석 기간: 2025.07.28~08.03",
    0.6, 6.48, 12.1, 0.65, size=13, color=C_NAVY, align=PP_ALIGN.CENTER)


# ════════════════════════════════════════════════════════════════════
# Slide 5: 연구지역 및 데이터
# ════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
add_rect(sl, 0, 0, 13.33, 7.5, fill=C_LGRAY)
slide_header(sl, "연구지역 및 데이터", "서울 성동구 | 2025년 여름 폭염 기간")

# 왼쪽: 역 목록
add_rect(sl, 0.4, 1.4, 4.5, 5.5, fill=C_WHITE)
add_rect(sl, 0.4, 1.4, 0.1, 5.5, fill=C_NAVY)
add_text(sl, "분석 대상 역 (7개)", 0.65, 1.5, 4.1, 0.45,
         size=15, bold=True, color=C_NAVY)

stations = [
    ("왕십리역", "2호선/수분당/중앙", "#E53935"),
    ("행당역",   "5호선",             "#FB8C00"),
    ("응봉역",   "경의중앙선",        "#8E24AA"),
    ("뚝섬역",   "2호선",             "#43A047"),
    ("성수역",   "2호선",             "#1E88E5"),
    ("서울숲역", "수인분당선",        "#00ACC1"),
    ("옥수역",   "3호선/중앙선",      "#6D4C41"),
]
for i, (name, line, _) in enumerate(stations):
    add_text(sl, f"● {name}", 0.65, 2.05 + i*0.56, 2.2, 0.5, size=13, bold=True, color=C_GRAY)
    add_text(sl, line, 2.85, 2.05 + i*0.56, 2.0, 0.5, size=11, color=RGBColor(0x75, 0x75, 0x75))

# 오른쪽: 데이터 표
add_rect(sl, 5.2, 1.4, 7.7, 5.5, fill=C_WHITE)
add_rect(sl, 5.2, 1.4, 7.7, 0.5, fill=C_BLUE)
add_text(sl, "사용 데이터", 5.4, 1.48, 7.3, 0.4, size=15, bold=True, color=C_WHITE)

data_rows = [
    ("S-DoT IoT 센서",  "성동구 57개",  "기온·습도·풍속 (1시간 단위)"),
    ("ASOS 기상관측",    "서울(108)",    "풍속·일사량 (2025.07.28~08.03)"),
    ("Open-Meteo",       "ERA5 기반",    "기온·습도·풍속·일사량 무료 API"),
    ("OSM 보행 네트워크","~15,000 링크", "성동구 전체 보행 경로"),
    ("건물 폴리곤",      "국토부 제공",  "층수 속성 → 높이 추정(층×3m)"),
    ("가로수 데이터",    "서울시 녹지",  "캐노피 비율 계산용"),
]
for i, (src, scope, desc) in enumerate(data_rows):
    bg = C_WHITE if i % 2 == 0 else C_LGRAY
    add_rect(sl, 5.2, 1.95 + i*0.75, 7.7, 0.72, fill=bg)
    add_text(sl, src,   5.35, 2.0 + i*0.75, 2.3, 0.6, size=12, bold=True,  color=C_NAVY)
    add_text(sl, scope, 7.7,  2.0 + i*0.75, 1.5, 0.6, size=11, color=C_BLUE)
    add_text(sl, desc,  9.25, 2.0 + i*0.75, 3.4, 0.6, size=11, color=C_GRAY)


# ════════════════════════════════════════════════════════════════════
# Slide 6: UTCI 계산 방법
# ════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
add_rect(sl, 0, 0, 13.33, 7.5, fill=C_LGRAY)
slide_header(sl, "UTCI 계산: SOLWEIG 기반 MRT → UTCI",
             "Lindberg et al. (2008) 방법론 적용")

# 수식 박스
add_rect(sl, 0.4, 1.4, 12.5, 2.5, fill=C_WHITE)
add_rect(sl, 0.4, 1.4, 0.12, 2.5, fill=C_ACCENT)
add_text(sl, "MRT 계산 (낮 시간)", 0.65, 1.5, 8.0, 0.45,
         size=15, bold=True, color=C_ACCENT)
add_text(sl,
    "MRT  =  Tair  +  0.5 × √(GHI × cos_z)  ×  SVF",
    0.65, 2.0, 11.8, 0.6, size=20, bold=True, color=C_NAVY, align=PP_ALIGN.CENTER)
add_text(sl,
    "SVF=1 (개활지): 태양복사 최대 → MRT 최고    |    SVF=0 (협곡): 복사 차단 → MRT ≈ Tair",
    0.65, 2.6, 11.8, 0.45, size=12, color=C_GRAY, align=PP_ALIGN.CENTER, italic=True)
add_text(sl,
    "UTCI  =  f (Tair, MRT, va, RH)    [pythermalcomfort 라이브러리, Bröde et al. 2012]",
    0.65, 3.1, 11.8, 0.55, size=17, bold=True, color=C_BLUE, align=PP_ALIGN.CENTER)

# 하단 3칸
cols3 = [
    (C_NAVY,   "기상 입력",
     "Open-Meteo Archive\n2025.07.28~08.03 7일 평균\n기온·습도·풍속·일사량(GHI)"),
    (C_BLUE,   "SVF 계산 (Oke 1987)",
     "SVF = 1/√(1+(H/W)²)\nH: 버퍼 20m 건물 평균 높이\nW: 도로 유형별 표준 폭"),
    (C_GREEN,  "캐노피 보정 (Chen & Ng 2012)",
     "ΔUTCI = 2.5°C × 캐노피비율\n가로수 버퍼 15m 면적 비율\n태양고도 가중치 적용"),
]
for i, (col, title, desc) in enumerate(cols3):
    x0 = 0.4 + i * 4.3
    add_rect(sl, x0, 4.0, 4.1, 3.1, fill=C_WHITE)
    add_rect(sl, x0, 4.0, 4.1, 0.5, fill=col)
    add_text(sl, title, x0+0.15, 4.07, 3.8, 0.4, size=13, bold=True, color=C_WHITE)
    add_text(sl, desc,  x0+0.2,  4.62, 3.7, 2.2, size=12.5, color=C_GRAY)


# ════════════════════════════════════════════════════════════════════
# Slide 7: 링크 회피 모델
# ════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
add_rect(sl, 0, 0, 13.33, 7.5, fill=C_LGRAY)
slide_header(sl, "링크 회피 모델 (Hard-Cut)", "UTCI ≥ 38°C 링크를 그래프에서 제거")

# 좌: 개념 설명
add_rect(sl, 0.4, 1.4, 5.8, 5.5, fill=C_WHITE)
add_rect(sl, 0.4, 1.4, 0.12, 5.5, fill=C_RED)
add_text(sl, "모델 설계 근거", 0.65, 1.5, 5.3, 0.45, size=15, bold=True, color=C_RED)

items = [
    "UTCI ≥ 38°C = 매우 강한 열스트레스",
    "  → 보행자 자발적 경로 회피",
    "  (Bröde et al. 2012 국제 기준)",
    "",
    "Classic PPA: 전체 그래프",
    "  → 15분 내 도달 가능 노드 집합",
    "",
    "Thermal PPA: 고온 링크 제거 후",
    "  → 우회로가 없으면 접근 불가",
    "",
    "감소율 = (Classic - Thermal)",
    "         ─────────────────── × 100",
    "              Classic",
]
for i, item in enumerate(items):
    if item == "":
        continue
    indent = item.startswith("  ")
    add_text(sl, ("▪ " if not indent and not item.startswith("─") else "") + item.strip(),
             0.65 + (0.3 if indent else 0), 2.0 + i*0.4,
             5.2 - (0.3 if indent else 0), 0.38,
             size=12.5 if indent else 13, color=C_GRAY,
             bold=(not indent and not item.startswith("─") and not item.startswith("감")))

# 우: UTCI 카테고리 표
add_rect(sl, 6.7, 1.4, 6.2, 5.5, fill=C_WHITE)
add_rect(sl, 6.7, 1.4, 6.2, 0.5, fill=C_NAVY)
add_text(sl, "UTCI 열스트레스 카테고리 (Bröde et al. 2012)",
         6.9, 1.47, 5.8, 0.4, size=13, bold=True, color=C_WHITE)

categories = [
    ("< 26°C",    "No thermal stress",         RGBColor(0xE3,0xF2,0xFD), False),
    ("26~32°C",   "Moderate heat stress",       RGBColor(0xFF,0xF9,0xC4), False),
    ("32~38°C",   "Strong heat stress",         RGBColor(0xFF,0xE0,0xB2), False),
    ("38~46°C ★", "Very strong heat stress",    RGBColor(0xFF,0xCC,0xBC), True),
    ("> 46°C",    "Extreme heat stress",        RGBColor(0xFF,0xCD,0xD2), False),
]
for i, (rng, label, bg, highlight) in enumerate(categories):
    add_rect(sl, 6.7, 2.0 + i*0.78, 6.2, 0.74, fill=bg)
    if highlight:
        add_rect(sl, 6.7, 2.0 + i*0.78, 0.12, 0.74, fill=C_RED)
    add_text(sl, rng,   6.9, 2.07 + i*0.78, 1.5, 0.58,
             size=13, bold=highlight, color=C_RED if highlight else C_NAVY)
    add_text(sl, label, 8.45, 2.07 + i*0.78, 4.3, 0.58,
             size=12.5, bold=highlight, color=C_RED if highlight else C_GRAY)


# ════════════════════════════════════════════════════════════════════
# Slide 8: 주요 결과 – 히트맵
# ════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
add_rect(sl, 0, 0, 13.33, 7.5, fill=C_LGRAY)
slide_header(sl, "주요 결과: 역별 접근성 감소율",
             "SOLWEIG 기반 UTCI | UTCI ≥ 38°C 링크 회피 | 시간예산 15분")

add_img(sl, os.path.join(FIG_DIR, 'fig1_heatmap_comparison.png'),
        0.35, 1.35, 12.6, 5.0)

add_rect(sl, 0.35, 6.45, 12.6, 0.8, fill=RGBColor(0xE3, 0xF2, 0xFD))
add_text(sl,
    "두 방법 모두 동일한 취약성 순서: 응봉역 · 서울숲역 (>65%) > 왕십리역 · 뚝섬역 > 행당역 > 옥수역 > 성수역",
    0.55, 6.55, 12.1, 0.6, size=13, bold=True, color=C_NAVY, align=PP_ALIGN.CENTER)


# ════════════════════════════════════════════════════════════════════
# Slide 9: 시간대별 변화
# ════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
add_rect(sl, 0, 0, 13.33, 7.5, fill=C_LGRAY)
slide_header(sl, "시간대별 접근성 감소율 변화",
             "오전은 안전 → 낮 13~16시 급격한 접근성 상실")

add_img(sl, os.path.join(FIG_DIR, 'fig3_timeseries.png'),
        0.35, 1.35, 12.6, 4.9)

add_rect(sl, 0.35, 6.35, 12.6, 0.9, fill=RGBColor(0xFF, 0xF8, 0xE1))
bullet(sl, [
    "7시·10시: 제거 링크 없음 → 접근성 감소 0% (기온 낮아 임계값 미달)",
    "13시: 고온 링크 23.5% 제거 → 역에 따라 최대 69% 접근성 상실",
    "응봉역·서울숲역: 16시에도 감소율 유지 (교량·개방공간 밀집 → 냉각 지연)",
], 0.55, 6.42, 12.1, h_each=0.27, size=12.5, color=C_NAVY, marker="▶ ")


# ════════════════════════════════════════════════════════════════════
# Slide 10: 응봉역 vs 성수역 대조
# ════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
add_rect(sl, 0, 0, 13.33, 7.5, fill=C_LGRAY)
slide_header(sl, "응봉역 vs 성수역: 공간적 대조",
             "열취약 교통 결절 vs 녹지 인접 결절 (13시)")

add_img(sl, os.path.join(FIG_DIR, 'fig2_eungbong_vs_seongsu.png'),
        0.35, 1.3, 12.6, 5.15)

# 하단 해석
add_rect(sl, 0.35, 6.55, 6.0, 0.8, fill=RGBColor(0xFC, 0xE4, 0xEC))
add_rect(sl, 0.35, 6.55, 0.1, 0.8, fill=C_RED)
add_text(sl, "응봉역 (68.6% 감소): 중랑천 교량 집중, 높은 SVF → 그늘 없음",
         0.55, 6.63, 5.6, 0.6, size=12, color=C_NAVY)

add_rect(sl, 7.0, 6.55, 6.0, 0.8, fill=RGBColor(0xE8, 0xF5, 0xE9))
add_rect(sl, 7.0, 6.55, 0.1, 0.8, fill=C_GREEN)
add_text(sl, "성수역 (13.0% 감소): 밀집 건물 협곡 → SVF 낮음 → 자연 차열",
         7.2, 6.63, 5.6, 0.6, size=12, color=C_NAVY)


# ════════════════════════════════════════════════════════════════════
# Slide 11: 회귀분석 결과
# ════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
add_rect(sl, 0, 0, 13.33, 7.5, fill=C_LGRAY)
slide_header(sl, "회귀분석: 접근성 감소 결정 요인",
             "Option A (역 단위, n=7) + Option B (링크 단위, n=8,003)")

# 좌: Option A
add_rect(sl, 0.4, 1.4, 5.8, 5.5, fill=C_WHITE)
add_rect(sl, 0.4, 1.4, 0.12, 5.5, fill=C_BLUE)
add_text(sl, "Option A — 역 단위 Pearson 상관", 0.65, 1.5, 5.3, 0.45,
         size=14, bold=True, color=C_BLUE)

corr_data = [
    ("고온링크 비율",   "r = +0.783",  "p = 0.037 **", True),
    ("교량 링크 비율",  "r = +0.646",  "p = 0.117",    False),
    ("평균 SVF",       "r = +0.611",  "p = 0.145",    False),
    ("하천까지 거리",   "r = −0.322",  "p = 0.481",    False),
    ("평균 캐노피",     "r = +0.087",  "p = 0.853",    False),
]
for i, (var, r_val, p_val, sig) in enumerate(corr_data):
    bg = RGBColor(0xE3, 0xF2, 0xFD) if sig else C_WHITE
    add_rect(sl, 0.4, 2.05 + i*0.75, 5.8, 0.7, fill=bg)
    add_text(sl, var,   0.6,  2.12 + i*0.75, 2.4, 0.56, size=12.5, bold=sig, color=C_NAVY)
    add_text(sl, r_val, 3.05, 2.12 + i*0.75, 1.5, 0.56, size=12.5, bold=sig,
             color=C_BLUE if sig else C_GRAY)
    add_text(sl, p_val, 4.55, 2.12 + i*0.75, 1.5, 0.56, size=12,
             color=C_RED if sig else C_GRAY)

add_text(sl, "* n=7, 통계적 해석에 주의 필요",
         0.6, 5.8, 5.5, 0.4, size=10.5, italic=True, color=RGBColor(0x9E,0x9E,0x9E))

# 우: Option B
add_rect(sl, 6.7, 1.4, 6.2, 5.5, fill=C_WHITE)
add_rect(sl, 6.7, 1.4, 0.12, 5.5, fill=C_ORANGE)
add_text(sl, "Option B — 링크 단위 로지스틱 회귀", 6.95, 1.5, 5.7, 0.45,
         size=14, bold=True, color=C_ORANGE)
add_text(sl, "DV: is_hot (UTCI ≥ 38°C = 1)  |  Pseudo R² = 0.951",
         6.95, 1.95, 5.7, 0.38, size=11.5, italic=True, color=C_GRAY)

logit_data = [
    ("건물 평균 높이↑", "OR = 0.028 ***", "높을수록 협곡 → 열위험 ↓"),
    ("도로 폭↑",        "OR = 12.5 ***",  "넓을수록 개방 → 열위험 ↑"),
    ("교량 여부",       "OR = 0.363 n.s.","높이·폭 통제 후 비유의미"),
]
for i, (var, or_val, interp) in enumerate(logit_data):
    is_sig = "***" in or_val
    add_rect(sl, 6.7, 2.45 + i*1.3, 6.2, 1.2,
             fill=RGBColor(0xFF,0xF3,0xE0) if is_sig else C_WHITE)
    add_rect(sl, 6.7, 2.45 + i*1.3, 0.12, 1.2,
             fill=C_ORANGE if is_sig else RGBColor(0xBD,0xBD,0xBD))
    add_text(sl, var,    6.95, 2.52 + i*1.3, 5.7, 0.45, size=13, bold=True,  color=C_NAVY)
    add_text(sl, or_val, 6.95, 2.95 + i*1.3, 5.7, 0.38, size=13, bold=is_sig, color=C_ORANGE if is_sig else C_GRAY)
    add_text(sl, "→ " + interp, 6.95, 3.3 + i*1.3, 5.7, 0.35,
             size=11.5, italic=True, color=C_GRAY)


# ════════════════════════════════════════════════════════════════════
# Slide 12: 결론
# ════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
add_rect(sl, 0, 0, 13.33, 7.5, fill=C_LGRAY)
slide_header(sl, "결론 및 시사점")

findings = [
    (C_RED,    "발견 1",
     "폭염 시 최대 69% 접근성 상실",
     "13시 기준 응봉역 68.6%, 서울숲역 67.8% — 10명 중 7명이 접근 범위 축소"),
    (C_BLUE,   "발견 2",
     "취약성은 도시 형태(H/W)로 결정",
     "건물 높이↑·도로 폭↓ → SVF↓ → 그늘 효과 → 열스트레스 위험 감소 (OR=0.028***)"),
    (C_GREEN,  "발견 3",
     "교량·수변 구간이 열 핫스팟",
     "응봉역: 중랑천 교량 밀집 → SVF 높음 → 그늘 없음 → 접근성 대폭 감소"),
    (C_ORANGE, "정책 시사점",
     "도시 설계 개입의 우선순위 제시",
     "교량·광장에 차양·녹화 우선 투자 | 15분 접근성 평가에 열환경 지표 포함 필요"),
]
for i, (col, tag, title, desc) in enumerate(findings):
    y0 = 1.4 + i * 1.45
    add_rect(sl, 0.4, y0, 12.5, 1.35, fill=C_WHITE)
    add_rect(sl, 0.4, y0, 0.12, 1.35, fill=col)
    add_rect(sl, 0.52, y0 + 0.1, 1.1, 0.42, fill=col)
    add_text(sl, tag,   0.56, y0 + 0.12, 1.0, 0.38, size=11, bold=True, color=C_WHITE)
    add_text(sl, title, 1.75, y0 + 0.1,  9.5, 0.45, size=15, bold=True, color=col)
    add_text(sl, desc,  1.75, y0 + 0.6,  10.8, 0.65, size=12.5, color=C_GRAY)


# ════════════════════════════════════════════════════════════════════
# Slide 13: 한계 및 향후 연구
# ════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
add_rect(sl, 0, 0, 13.33, 7.5, fill=C_LGRAY)
slide_header(sl, "한계 및 향후 연구")

# 한계
add_rect(sl, 0.4, 1.4, 5.8, 5.5, fill=C_WHITE)
add_rect(sl, 0.4, 1.4, 0.12, 5.5, fill=C_RED)
add_text(sl, "연구의 한계", 0.65, 1.5, 5.3, 0.4, size=15, bold=True, color=C_RED)
limits = [
    "역 기준 Catchment의 한계",
    "  개인 이동 기점·목적지 미반영",
    "  실제 경로 선택 행동 단순화",
    "합성 DSM의 정확도 제한",
    "  LiDAR 미확보 → 층수×3m 추정",
    "단일 기상 지점 입력",
    "  성동구 전체 균일 기상 적용",
    "  미기후 공간 변동성 과소 추정",
    "역 7개로 통계적 검정력 제한",
    "  회귀 결과는 탐색적 수준",
]
for i, item in enumerate(limits):
    indent = item.startswith("  ")
    add_text(sl, ("▪ " if not indent else "") + item.strip(),
             0.65 + (0.25 if indent else 0), 2.0 + i*0.43,
             5.2, 0.4, size=12 if indent else 13,
             color=RGBColor(0x75,0x75,0x75) if indent else C_GRAY)

# 향후 — STP 확장 강조
add_rect(sl, 6.7, 1.4, 6.2, 5.5, fill=C_WHITE)
add_rect(sl, 6.7, 1.4, 0.12, 5.5, fill=C_BLUE)
add_text(sl, "향후 연구 방향", 6.95, 1.5, 5.7, 0.4, size=15, bold=True, color=C_BLUE)

# STP 확장 강조 박스
add_rect(sl, 6.82, 2.0, 5.95, 1.5, fill=RGBColor(0xE3,0xF2,0xFD))
add_rect(sl, 6.82, 2.0, 0.1, 1.5, fill=C_BLUE)
add_text(sl, "★  Space-Time Prism(STP)으로 확장",
         7.05, 2.08, 5.6, 0.45, size=13, bold=True, color=C_NAVY)
add_text(sl, "개인 이동 기점·목적지 + 시간예산 동시 반영\n→ 폭염의 시공간 이동 제약을 개인 단위로 정량화",
         7.05, 2.52, 5.6, 0.85, size=12, color=C_GRAY)

futures = [
    "LiDAR 기반 정밀 DSM 확보",
    "  → UTCI 공간 정확도 향상",
    "다중 기상 관측 포인트 적용",
    "  → 미기후 공간 변동 반영",
    "서울시 전역으로 확장",
    "  → 자치구별 열 취약성 비교",
    "그늘막·가로수 설치 효과 시뮬레이션",
]
for i, item in enumerate(futures):
    indent = item.startswith("  ")
    add_text(sl, ("▪ " if not indent else "") + item.strip(),
             6.95 + (0.25 if indent else 0), 3.65 + i*0.43,
             5.7, 0.4, size=12 if indent else 13,
             color=RGBColor(0x75,0x75,0x75) if indent else C_GRAY)


# ════════════════════════════════════════════════════════════════════
# Slide 14: 감사합니다
# ════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
add_rect(sl, 0, 0, 13.33, 7.5, fill=C_NAVY)
add_rect(sl, 0, 3.5, 13.33, 0.06, fill=C_BLUE)

add_text(sl, "감사합니다", 0, 1.5, 13.33, 1.5,
         size=52, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
add_text(sl, "Thank you for your attention",
         0, 3.0, 13.33, 0.7, size=20, color=RGBColor(0x90,0xCA,0xF9),
         align=PP_ALIGN.CENTER, italic=True)

add_text(sl, "신 진  |  경희대학교 기후사회과학융합학과  |  도시컴퓨팅연구실",
         0, 4.2, 13.33, 0.6, size=16, color=C_WHITE, align=PP_ALIGN.CENTER)

# 참고문헌 간략
refs = [
    "Bröde et al. (2012) Int. J. Biometeorology · Colaninno et al. (2024) Urban Analytics & City Science",
    "Lindberg et al. (2008, 2018) · Oke (1987) · Moreno et al. (2021) Smart Cities",
    "Hägerstrand (1970) · Miller (1991) · Chen & Ng (2012) · Hersbach et al. (2020)",
]
add_rect(sl, 1.5, 5.2, 10.33, 1.9, fill=RGBColor(0x0D, 0x47, 0xA1))
add_text(sl, "주요 참고문헌", 1.7, 5.28, 9.9, 0.4,
         size=12, bold=True, color=RGBColor(0x90,0xCA,0xF9))
for i, ref in enumerate(refs):
    add_text(sl, ref, 1.7, 5.68 + i*0.42, 9.9, 0.4,
             size=10.5, color=RGBColor(0xBB,0xDE,0xFB))


# ── 저장 ──────────────────────────────────────────────────────────────
os.makedirs(BASE, exist_ok=True)
prs.save(OUT)
print(f"PPT 저장 완료: {OUT}")
print(f"슬라이드 수: {len(prs.slides)}장")
