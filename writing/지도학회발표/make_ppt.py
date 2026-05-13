"""
TAVI 학회 발표 PPT 생성
스타일: Traffic-IT_발제_신진.pptx 동일
- 슬라이드 텍스트 충실 작성
- 슬라이드 노트(대본) 포함
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from PIL import Image
import os

# ── 색상 ─────────────────────────────────────────────────────────────────
C_RED     = RGBColor(0xB2, 0x22, 0x34)
C_DARKRED = RGBColor(0x8B, 0x00, 0x00)
C_WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
C_BLACK   = RGBColor(0x1A, 0x1A, 0x1A)
C_GRAY    = RGBColor(0x55, 0x55, 0x55)
FONT      = '맑은 고딕'

FIG_DIR  = '/Users/jin/석사논문/TAVI/Thermal_Catchment/03_결과물/figures'
PROF_DIR = os.path.join(FIG_DIR, '교수님논의')
OUT_PATH = '/Users/jin/석사논문/TAVI/writing/지도학회발표/2026-05-13_TAVI_학회발표.pptx'

# ── 헬퍼 ─────────────────────────────────────────────────────────────────
def add_textbox(slide, text, left, top, width, height,
                size=12, bold=False, color=C_BLACK, align=PP_ALIGN.LEFT):
    txb = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = FONT
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return txb

def add_multiline(slide, lines, left, top, width, height, size=12, color=C_BLACK):
    """lines: list of (text, bold)"""
    txb = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txb.text_frame
    tf.word_wrap = True
    for i, (text, bold) in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        run = p.add_run()
        run.text = text
        run.font.name = FONT
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = color
    return txb

def add_rect(slide, left, top, width, height, fill_color,
             text=None, text_size=12, text_bold=True, text_color=C_WHITE):
    shape = slide.shapes.add_shape(
        1, Inches(left), Inches(top), Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    if text:
        tf = shape.text_frame
        tf.word_wrap = False
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = text
        run.font.name = FONT
        run.font.size = Pt(text_size)
        run.font.bold = text_bold
        run.font.color.rgb = text_color
    return shape

def add_divider(slide, top=1.18):
    add_rect(slide, 0.40, top, 12.53, 0.02, C_RED)

def add_chapter_header(slide, chapter_num, chapter_title, section_label=None):
    add_textbox(slide, f'Chapter {chapter_num}.', 0.40, 0.18, 4.0, 0.40,
                size=11, color=C_RED)
    add_textbox(slide, chapter_title, 0.40, 0.52, 12.50, 0.55,
                size=24, color=C_BLACK)
    add_divider(slide, top=1.18)
    if section_label:
        label_w = max(len(section_label) * 0.19 + 0.5, 2.0)
        add_rect(slide, 0.40, 1.35, label_w, 0.36,
                 C_DARKRED, text=section_label, text_size=12)

def add_section_box(slide, label, top):
    label_w = max(len(label) * 0.19 + 0.5, 2.0)
    add_rect(slide, 0.40, top, label_w, 0.36,
             C_DARKRED, text=label, text_size=12)

def add_image_fit(slide, img_path, left, top, max_w, max_h,
                  caption=None, caption_size=10):
    img = Image.open(img_path)
    iw, ih = img.size
    ratio = min(max_w / iw, max_h / ih)
    w, h = iw * ratio, ih * ratio
    offset_x = (max_w - w) / 2
    slide.shapes.add_picture(
        img_path,
        Inches(left + offset_x), Inches(top),
        Inches(w), Inches(h)
    )
    if caption:
        add_textbox(slide, caption,
                    left, top + h + 0.05, max_w, 0.35,
                    size=caption_size, color=C_GRAY, align=PP_ALIGN.CENTER)
    return w, h

def set_notes(slide, text):
    notes_slide = slide.notes_slide
    tf = notes_slide.notes_text_frame
    tf.text = text


# ── PPT 초기화 ────────────────────────────────────────────────────────────
prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.50)
blank = prs.slide_layouts[6]

def new_slide():
    return prs.slides.add_slide(blank)


# ═══════════════════════════════════════════════════════════════════════════
# 슬라이드 1: 표지
# ═══════════════════════════════════════════════════════════════════════════
s = new_slide()

add_textbox(s, '한국지도학회 학술대회 발표', 0.40, 0.20, 8.0, 0.40,
            size=11, color=C_GRAY)
add_textbox(s,
    'Thermal Catchment Area:\nMRT 기반 확률적 임계값과 Monte Carlo 접근',
    0.80, 0.90, 8.20, 2.20, size=28, bold=True, color=C_BLACK)
add_textbox(s,
    '폭염을 반영한 보행 대중교통 접근성 공간 단위 제안 — 서울 성동구 사례',
    0.80, 3.20, 8.20, 0.60, size=14, color=C_GRAY)
add_textbox(s, '신진', 9.50, 5.10, 3.50, 0.40,
            size=15, bold=True, color=C_BLACK)
add_textbox(s, '경희대학교 기후사회과학융합학과 석사과정', 9.50, 5.55, 3.50, 0.35,
            size=10, color=C_BLACK)
add_textbox(s, '2026. 05. 30', 9.50, 5.95, 3.50, 0.35,
            size=11, color=C_BLACK)
add_rect(s, 0.40, 6.85, 12.53, 0.02, C_RED)

img_path = os.path.join(FIG_DIR, 'tarr_spatial_map.png')
if os.path.exists(img_path):
    add_image_fit(s, img_path, 9.00, 0.50, 4.10, 4.40)

set_notes(s,
"""안녕하세요. 경희대학교 기후사회과학융합학과 석사과정 신진입니다.
오늘 발표할 주제는 'Thermal Catchment Area'입니다.

폭염이 심해지는 여름철, 우리가 지하철역에 걸어서 갈 수 있는 범위가 열환경에 의해 실제로 줄어든다는 점에 착안하여,
MRT(평균복사온도)를 기반으로 한 새로운 보행 접근성 공간 단위를 제안하고,
그 임계값의 불확실성을 Monte Carlo 방법으로 정량화한 연구를 소개드리겠습니다.
발표 시간은 약 15분입니다.""")


# ═══════════════════════════════════════════════════════════════════════════
# 슬라이드 2: 목차
# ═══════════════════════════════════════════════════════════════════════════
s = new_slide()
add_textbox(s, '목차', 0.40, 0.18, 4.0, 0.55, size=24, color=C_BLACK)
add_divider(s, top=0.82)

chapters = [
    ('Chapter 1.', '연구 배경 및 목적'),
    ('Chapter 2.', '연구 지역 및 데이터'),
    ('Chapter 3.', '방법론 — MRT 산출 · Thermal Catchment · Monte Carlo'),
    ('Chapter 4.', '분석 결과 — S-curve · 공간 분포 · 시간대 비교 · 회귀분석'),
    ('Chapter 5.', '한계점 및 향후 연구'),
    ('Chapter 6.', '결론'),
]
for i, (ch, title) in enumerate(chapters):
    row_top = 1.10 + i * 0.95
    add_rect(s, 0.40, row_top, 0.08, 0.36, C_RED)
    add_textbox(s, ch,    0.60, row_top, 2.20, 0.36, size=12, bold=True,  color=C_RED)
    add_textbox(s, title, 2.80, row_top, 9.80, 0.36, size=13, bold=False, color=C_BLACK)

set_notes(s,
"""목차를 간략히 소개드리겠습니다.
먼저 연구 배경과 목적, 연구 지역과 데이터를 설명드리고,
방법론 파트에서는 MRT 산출 방법, Thermal Catchment 개념, 그리고 Monte Carlo 설계를 설명드립니다.
분석 결과에서는 S-curve 발견, TARR 공간 분포, 시간대별 비교, 회귀분석 순으로 보여드리겠습니다.
마지막으로 한계점과 결론으로 마무리하겠습니다.""")


# ═══════════════════════════════════════════════════════════════════════════
# 슬라이드 3: Chapter 1 — 연구 배경 (1)
# ═══════════════════════════════════════════════════════════════════════════
s = new_slide()
add_chapter_header(s, 1, '연구 배경', '기존 접근성 분석의 한계')

add_multiline(s, [
    ('● 역세권 Catchment Area는 보행 네트워크의 물리적 거리·시간만을 기준으로 정의됩니다.', True),
    ('   → 보행자가 균질한 환경에서 이동한다는 암묵적 가정에 기반합니다.', False),
    ('   → 도로 유형, 경사, 기상 조건 등 실제 보행 환경의 이질성을 전혀 반영하지 못합니다.', False),
    ('', False),
    ('● 기후변화로 인해 도시 폭염의 빈도와 강도가 지속적으로 증가하고 있습니다.', True),
    ('   → 서울 기준 폭염일수: 2010년대 평균 10.3일 → 2020년대 평균 15.8일 (기상청)', False),
    ('   → 2025년 서울 성동구: 7월 28일~8월 3일 연속 7일 폭염 지속', False),
    ('', False),
    ('● UTCI(Universal Thermal Climate Index) ≥ 38°C는 "Very Strong Heat Stress" 단계입니다.', True),
    ('   → 이 조건에서는 보행 활동이 신체에 심각한 위험을 초래하며,', False),
    ('   → 실제로 보행자들은 노출 경로를 회피하거나 이동 자체를 포기합니다.', False),
], 0.45, 1.88, 12.40, 4.20, size=12)

add_section_box(s, '연구 공백', 6.20)
add_multiline(s, [
    ('● 열환경 임계값을 초과할 때 접근 가능 범위 자체가 달라지는 공간 단위 개념은 제안된 바 없습니다.', False),
    ('● 임계값의 불확실성(기상 조건 변동)을 확률적으로 다루는 연구도 존재하지 않습니다.', False),
], 0.45, 6.70, 12.40, 0.70, size=12)

set_notes(s,
"""우리가 흔히 역세권 분석에서 사용하는 Catchment Area, 즉 역에서 걸어서 갈 수 있는 범위는
물리적 거리나 시간만을 기준으로 합니다.
그런데 폭염이 심한 여름철에는 어떨까요?
사람들이 뜨거운 햇볕이 내리쬐는 길을 피해 다른 경로를 선택하거나, 아예 포기하는 경우가 생깁니다.
UTCI 38°C 이상은 국제 기준으로 "매우 강한 열 스트레스" 단계로, 이 조건에서 보행은 건강에 심각한 위험입니다.
그런데 기존 연구 어디에도 이 열환경 제약을 공간 단위 자체에 반영한 연구는 없습니다.
이것이 바로 저희 연구가 출발한 지점입니다.""")


# ═══════════════════════════════════════════════════════════════════════════
# 슬라이드 4: Chapter 1 — 연구 목적 & 질문
# ═══════════════════════════════════════════════════════════════════════════
s = new_slide()
add_chapter_header(s, 1, '연구 목적 및 질문', '연구 목적')

add_multiline(s, [
    ('● Thermal Catchment Area 개념 제안', True),
    ('   → 폭염 시 MRT(평균복사온도) 기반 Hard Cut을 적용하여', False),
    ('      실제 보행 가능 범위를 기존 Catchment와 구분되는 새로운 공간 단위로 정식화합니다.', False),
    ('', False),
    ('● 서울 성동구 집계구 기점 대중교통 접근성 분석에 적용하여 실증적 유효성을 검증합니다.', True),
    ('   → 집계구(n=506) × 정류장(489개, 네트워크 노드 385개) 기반 Dijkstra 최단경로 분석', False),
    ('', False),
    ('● Monte Carlo 방법으로 임계값 불확실성을 정량화합니다.', True),
    ('   → threshold ~ N(μ=55°C, σ²=4²) 에서 2,000개 샘플링 → TARR 확률 분포 획득', False),
], 0.45, 1.88, 12.40, 3.20, size=12)

add_section_box(s, '연구 질문', 5.20)
add_multiline(s, [
    ('Q1. MRT Hard Cut 적용 시, 폭염 조건에서 기존 Catchment 대비 Thermal Catchment의 접근 가능 정류장은 얼마나 감소하는가?', True),
    ('', False),
    ('Q2. MRT 임계값의 불확실성을 Monte Carlo로 전파했을 때, TARR의 확률 분포와 공간 패턴은 어떠한가?', True),
    ('', False),
    ('Q3. TARR 패턴이 집계구의 공간 환경 변수(SVF, 캐노피, 건물 형태 등)와 어떠한 관계를 갖는가?', True),
], 0.45, 5.70, 12.40, 1.60, size=12)

set_notes(s,
"""연구 목적은 세 가지입니다.
첫째, Thermal Catchment Area라는 새로운 공간 단위 개념을 제안합니다.
둘째, 서울 성동구를 대상으로 실제 데이터로 검증합니다.
셋째, 임계값이 얼마냐에 따라 결과가 얼마나 달라지는지를 Monte Carlo로 정량화합니다.

연구 질문도 이 세 가지에 대응합니다.
첫 번째 질문은 폭염 시 실제로 정류장 접근이 얼마나 줄어드는가,
두 번째는 임계값이 불확실할 때 그 결과의 범위가 어느 정도인가,
세 번째는 어떤 공간 환경에서 접근성 손실이 더 크게 나타나는가입니다.""")


# ═══════════════════════════════════════════════════════════════════════════
# 슬라이드 5: Chapter 2 — 연구 지역 & 데이터
# ═══════════════════════════════════════════════════════════════════════════
s = new_slide()
add_chapter_header(s, 2, '연구 지역 및 데이터', '연구 지역: 서울특별시 성동구')

add_multiline(s, [
    ('● 면적: 16.86 km² | 지하철역 7개소: 왕십리·행당·응봉·뚝섬·성수·서울숲·옥수', True),
    ('● 한강·중랑천 수변 공간 + 서울숲 대형 녹지 + 왕십리·성수 업무지구 + 아파트 단지 혼재', False),
    ('   → 도시 형태가 다양하여 열환경 공간 변이가 크고, 접근성 분석의 적절한 대상 지역', False),
    ('', False),
], 0.45, 1.88, 7.50, 1.80, size=12)

add_section_box(s, '분석 데이터', 3.80)
add_multiline(s, [
    ('● 보행 네트워크:  OpenStreetMap (osmnx) — 보행 링크 15,608개', False),
    ('● 기상 데이터:     성동구 S-DoT 스마트 센서 — 57개소, 10분 간격 (기온·습도·풍속)', False),
    ('● 일사량(GHI):    Open-Meteo archive — 성동구 중심 단일 지점', False),
    ('● 집계구 경계:    통계청 2016 기준 — 570개 집계구 (유효 집계구 506개)', False),
    ('● 대중교통 정류장: 지하철 7개 + 버스 482개 = 489개 (네트워크 노드 기준 385개)', False),
    ('● 건물 데이터:    OSM 건물 폴리곤 — SVF 산출용', False),
    ('', False),
    ('● 분석 기간: 2025년 7월 28일 – 8월 3일 폭염일 7일 평균', True),
    ('● 분석 시간대: 13시 (태양 복사 최대·폭염 피크), 비교 분석: 9시·13시·18시', True),
], 0.45, 4.30, 7.50, 2.80, size=12)

img_path = os.path.join(PROF_DIR, '00_SDot센서위치.png')
if os.path.exists(img_path):
    add_image_fit(s, img_path, 8.20, 1.30, 4.90, 5.80,
                  caption='그림 1. 성동구 S-DoT 기상 센서 위치 (n=57개)')

set_notes(s,
"""연구 지역은 서울 성동구입니다.
성동구는 서울숲, 한강, 중랑천 등 녹지와 수변 공간이 있는 반면,
왕십리·성수 등 고밀도 업무지구와 아파트 단지도 혼재하여
열환경이 공간적으로 매우 다양합니다.
그래서 Thermal Catchment가 공간적으로 어떻게 달라지는지 보기에 최적의 연구 지역입니다.

오른쪽 그림은 분석에 활용한 S-DoT 기상 센서 57개소의 위치입니다.
성동구 전역에 약 300~500m 간격으로 설치되어 있고,
이 센서에서 기온, 상대습도, 풍속을 수집해 MRT 계산에 활용했습니다.
대중교통 정류장은 총 489개이며, 보행 네트워크 노드 기준으로는 385개를 사용했습니다.""")


# ═══════════════════════════════════════════════════════════════════════════
# 슬라이드 6: Chapter 3 — MRT 산출 방법론
# ═══════════════════════════════════════════════════════════════════════════
s = new_slide()
add_chapter_header(s, 3, '방법론', 'Step 1: MRT 산출 — SOLWEIG 약식 구현')

add_multiline(s, [
    ('● Lindberg & Grimmond(2011) SOLWEIG 모형을 기반으로 Python으로 직접 구현하였습니다.', True),
    ('● 보행 네트워크 링크(15,608개) 각각에 대해 MRT를 산출합니다.', True),
    ('', False),
    ('  Tmrt  =  [ (Ksw + Llw) / σ ]^0.25  − 273.15   [°C]', True),
    ('', False),
    ('  ▶ 단파복사 Ksw  =  αp × [ (1−SVF)·Kdif + shadow·Kdir·cosθ + SVF·Kdif ]', False),
    ('     - SVF (Sky View Factor): 하늘이 얼마나 열려 있는가 (0=완전 차폐, 1=완전 개방)', False),
    ('     - Kdir (직달일사), Kdif (확산일사): GHI → Erbs 모델로 분리', False),
    ('     - shadow: 13시 태양 천정각 기반 그림자 마스크 (1=그늘, 0=직사광선)', False),
    ('     - αp = 0.70 (인체 흡수율)', False),
    ('', False),
    ('  ▶ 장파복사 Llw  =  εp × [ SVF·Lsky + (1−SVF)·Lwall ]', False),
    ('     - Lsky = εsky·σ·Tair⁴  (Brutsaert 1975 대기 방출률)', False),
    ('     - Lwall = εw·σ·(Tair + 5)⁴  (건물벽면 복사, ΔT=5K 가정)', False),
    ('     - εp = 0.97 (인체 방출률)', False),
], 0.45, 1.88, 7.80, 5.20, size=11)

img_path = os.path.join(PROF_DIR, '06_MRT.png')
if os.path.exists(img_path):
    add_image_fit(s, img_path, 8.40, 1.50, 4.70, 5.60,
                  caption='그림 2. 링크별 MRT 분포 (13시, 폭염일 평균)')

set_notes(s,
"""방법론의 첫 번째 단계는 MRT 산출입니다.
MRT, 즉 평균복사온도는 보행자가 사방에서 받는 복사 에너지를 통합한 지표로,
열 스트레스 체감에 가장 큰 영향을 주는 요소입니다.

저희는 SOLWEIG 모형을 Python으로 직접 구현해서
보행 네트워크 링크 1만 5천여 개 각각에 대해 MRT를 산출했습니다.

가장 중요한 변수는 SVF입니다.
SVF가 낮을수록 건물이 하늘을 많이 가린다는 의미이고,
그만큼 직사광선도 차단되어 MRT가 낮아집니다.
오른쪽 그림을 보시면, 건물이 밀집한 구역의 링크들은 MRT가 상대적으로 낮고,
한강변이나 서울숲 주변 개방 구역은 MRT가 높게 나타나는 것을 확인할 수 있습니다.""")


# ═══════════════════════════════════════════════════════════════════════════
# 슬라이드 7: Chapter 3 — Thermal Catchment & TARR
# ═══════════════════════════════════════════════════════════════════════════
s = new_slide()
add_chapter_header(s, 3, '방법론', 'Step 2: Thermal Catchment Area & TARR 정의')

add_multiline(s, [
    ('● Classic Catchment', True),
    ('   → 집계구 도심에서 보행 속도 4.5 km/h 기준, 15분 이내 도달 가능한 정류장의 수', False),
    ('   → 열환경 무시 — 기존 역세권 접근성 분석의 표준 방식', False),
    ('', False),
    ('● Thermal Catchment  [본 연구 핵심 개념]', True),
    ('   → MRT ≥ 임계값(threshold)인 링크를 네트워크에서 제거한 뒤,', False),
    ('      남은 네트워크에서 15분 이내 도달 가능한 정류장의 수', False),
    ('   → Hard Cut 가정: 보행자는 MRT가 임계값 이상인 링크를 완전히 회피', False),
    ('   → Dijkstra 알고리즘(최단경로)으로 각 집계구에서 도달 가능 정류장 탐색', False),
    ('', False),
    ('● TARR (Thermal Accessibility Reduction Rate) — 열환경 접근성 감소율', True),
    ('', False),
    ('   TARR (%) =  ( Classic_cnt − Thermal_cnt ) / Classic_cnt  ×  100', True),
    ('', False),
    ('   → TARR = 0%:   열환경 제약 없음 — 폭염에도 접근 가능 범위 유지', False),
    ('   → TARR = 100%: 모든 정류장 접근 불가 — 완전 고립', False),
    ('   → TARR이 높은 집계구 = 폭염 시 대중교통 접근성 취약 지역', False),
], 0.45, 1.88, 12.40, 5.20, size=12)

set_notes(s,
"""두 번째 단계는 Thermal Catchment Area 개념 정의입니다.

기존의 Classic Catchment는 4.5 km/h 보행 속도로 15분 안에 갈 수 있는 정류장 수를 셉니다.
여기에 열환경 제약을 추가한 것이 Thermal Catchment입니다.
MRT가 임계값 이상인 링크를 네트워크에서 완전히 제거한 뒤,
남은 네트워크로만 15분 내 도달 가능한 정류장을 세는 방식입니다.

이때 "제거"하는 방식, 즉 Hard Cut은 보행자가 뜨거운 길을 완전히 피해 돌아간다는 가정입니다.

두 값의 차이를 기존 값으로 나누어 퍼센트로 표현한 것이 TARR입니다.
TARR이 클수록 폭염 시 접근성 손실이 크다는 의미로,
이 값이 높은 집계구가 바로 폭염 취약 지역입니다.""")


# ═══════════════════════════════════════════════════════════════════════════
# 슬라이드 8: Chapter 3 — Monte Carlo 설계
# ═══════════════════════════════════════════════════════════════════════════
s = new_slide()
add_chapter_header(s, 3, '방법론', 'Step 3: Monte Carlo 임계값 불확실성 분석')

add_multiline(s, [
    ('● 왜 Monte Carlo가 필요한가?', True),
    ('   → MRT 임계값을 "55°C"로 단일 고정하면 임계값 선택 자체의 불확실성이 무시됩니다.', False),
    ('   → 기상 조건(기온·습도·풍속)은 날마다 달라지므로, "위험 MRT" 역시 범위로 존재합니다.', False),
    ('   → 교수님 피드백: "임계값이 애매할 때 Monte Carlo를 돌려라"', False),
    ('', False),
    ('● 분석 설계 (N=2,000 샘플, 난수 시드 42 고정)', True),
    ('', False),
    ('   Step 1.  MRT 임계값 그리드 [45, 49, 51, 53, 55, 57, 59, 61, 65°C] 설정', False),
    ('            → 9개 임계값 각각에 대해 집계구별(n=506) TARR 사전 계산', False),
    ('', False),
    ('   Step 2.  threshold ~ N(μ=55°C, σ=4°C) 에서 2,000개 샘플링', False),
    ('            → 95% 범위: 47~63°C (물리적으로 가능한 MRT 임계값 범위)', False),
    ('', False),
    ('   Step 3.  집계구별 TARR(threshold) 선형 보간 → 샘플별 TARR 획득', False),
    ('            → (집계구 506 × 샘플 2,000) 매트릭스 구성', False),
    ('', False),
    ('   Step 4.  집계구별 TARR 중앙값·95% CI(2.5~97.5 퍼센타일) 산출', False),
], 0.45, 1.88, 12.40, 5.20, size=12)

set_notes(s,
"""세 번째 단계는 Monte Carlo를 이용한 임계값 불확실성 분석입니다.

55°C라는 임계값은 UTCI 38°C 역산 등으로 설정했지만, 이 값 자체가 확실한 것은 아닙니다.
기상 조건이 조금만 달라져도 "위험 MRT"는 변합니다.
그래서 임계값을 평균 55°C, 표준편차 4°C인 정규분포에서 2천 개 샘플링해서,
각 샘플마다 TARR을 계산하는 방식입니다.

구체적으로는 9개 임계값 지점에서 TARR을 먼저 계산해두고,
샘플링한 임계값에 대해 선형 보간으로 TARR을 추정합니다.
이렇게 하면 집계구마다 2천 개의 TARR 값이 생기고,
그 분포의 중앙값과 95% 신뢰구간을 결과로 제시합니다.""")


# ═══════════════════════════════════════════════════════════════════════════
# 슬라이드 9: Chapter 4 — S-curve (핵심 결과)
# ═══════════════════════════════════════════════════════════════════════════
s = new_slide()
add_chapter_header(s, 4, '분석 결과', '핵심 발견: MRT 임계값 민감도 S-curve')

img_path = os.path.join(FIG_DIR, 'mc_threshold_sensitivity.png')
if os.path.exists(img_path):
    add_image_fit(s, img_path, 0.50, 1.55, 8.20, 5.60,
                  caption='그림 3. MRT 임계값 민감도 곡선 (집계구 평균 TARR ± 1 std, n=506, 13시)')

add_multiline(s, [
    ('핵심 해석', True),
    ('', False),
    ('● 45°C 임계값', False),
    ('   TARR ≈ 95%', True),
    ('   거의 모든 링크 차단', False),
    ('', False),
    ('● 53°C 임계값', False),
    ('   TARR ≈ 88%', True),
    ('', False),
    ('● 55°C 임계값 ← 기준값', False),
    ('   TARR ≈ 73%', True),
    ('', False),
    ('● 57°C 임계값', False),
    ('   TARR ≈ 39%', True),
    ('', False),
    ('● 65°C 임계값', False),
    ('   TARR ≈ 5%', True),
    ('   거의 차단 없음', False),
    ('', False),
    ('★ 53~57°C 구간:',  True),
    ('  TARR 88%→39% 급변', True),
    ('  = S-curve 변곡점', True),
    ('', False),
    ('→ 55°C = 변곡점 중심,', False),
    ('  보수적 기준값', False),
], 9.00, 1.55, 4.10, 5.60, size=11)

set_notes(s,
"""이것이 이번 연구에서 발견한 가장 중요한 결과입니다: S-curve입니다.

S자 형태의 반응 곡선을 보면,
45°C에서는 TARR이 95%로 거의 모든 정류장에 접근 불가,
65°C에서는 5%로 거의 제약 없음인데,
그 중간인 53~57°C 구간에서 TARR이 88%에서 39%로 급격히 떨어집니다.
2°C 차이에 49%p 변동이 일어나는 것입니다.

이것이 S-curve의 변곡점이고, 55°C는 그 중심값입니다.

Monte Carlo에서 CI 폭이 넓게 나온 것도 바로 이 때문입니다.
CI가 넓다는 것은 결과가 불안정하다는 의미가 아니라,
이 민감 구간이 존재한다는 것, 즉 55°C 주변에서 임계값 선택이 결과에 결정적 영향을 미친다는 것을 보여줍니다.""")


# ═══════════════════════════════════════════════════════════════════════════
# 슬라이드 10: Chapter 4 — TARR 공간 분포
# ═══════════════════════════════════════════════════════════════════════════
s = new_slide()
add_chapter_header(s, 4, '분석 결과', 'TARR 공간 분포 (MRT 55°C 기준, 13시)')

img_path = os.path.join(FIG_DIR, 'tarr_spatial_map.png')
if os.path.exists(img_path):
    add_image_fit(s, img_path, 0.40, 1.50, 9.00, 5.65,
                  caption='그림 4. 집계구별 TARR 공간 분포 (n=506, 폭염일 평균)')

add_multiline(s, [
    ('주요 통계', True),
    ('', False),
    ('● 집계구 평균 TARR: 73.3%', True),
    ('   → 폭염 시 평균 73%의', False),
    ('     정류장 접근 기회 손실', False),
    ('', False),
    ('● TARR > 80%: 고취약 지역', False),
    ('   (짙은 빨강)', False),
    ('', False),
    ('● 공간 패턴', True),
    ('   → 한강변·서울숲 주변', False),
    ('     TARR 높음', False),
    ('     (개방 구간 = 고온 링크 多)', False),
    ('', False),
    ('   → 건물 밀집 구간', False),
    ('     TARR 상대적 낮음', False),
    ('     (SVF↓ = 그늘 多)', False),
    ('', False),
    ('● 공간 군집 패턴 존재', False),
    ('   → Moran\'s I 분석 필요', False),
], 9.60, 1.55, 3.50, 5.60, size=11)

set_notes(s,
"""이 지도는 집계구별 TARR의 공간 분포입니다.
색이 짙을수록 폭염 시 접근성 손실이 크다는 의미입니다.

성동구 전체 평균 TARR은 73.3%로,
폭염 피크 시간인 13시에 평균적으로 73%의 정류장 접근 기회가 사라집니다.

공간 패턴을 보면,
한강변이나 서울숲 주변 개방된 구간은 TARR이 높고,
건물이 밀집한 구간은 상대적으로 낮습니다.
이는 개방된 구간일수록 MRT가 높아 더 많은 링크가 차단되기 때문입니다.

또한 지도에서 같은 색의 집계구들이 공간적으로 뭉쳐있는 경향이 보이는데,
이는 공간 자기상관(spatial autocorrelation)이 존재한다는 것을 시사합니다.""")


# ═══════════════════════════════════════════════════════════════════════════
# 슬라이드 11: Chapter 4 — Monte Carlo 분포
# ═══════════════════════════════════════════════════════════════════════════
s = new_slide()
add_chapter_header(s, 4, '분석 결과', 'Monte Carlo TARR 분포 및 불확실성 공간 분포')

img_path = os.path.join(FIG_DIR, 'mc_tarr_distribution.png')
if os.path.exists(img_path):
    add_image_fit(s, img_path, 0.40, 1.50, 12.50, 4.40,
                  caption='그림 5. (좌) 집계구별 TARR 중앙값 분포 히스토그램 / (우) 95% CI 폭 분포 히스토그램')

img_path2 = os.path.join(FIG_DIR, 'mc_tarr_ci_map.png')
if os.path.exists(img_path2):
    pass  # 다음 슬라이드에서 단독으로

add_multiline(s, [
    ('● TARR 중앙값 평균: 70.2%  (std: 18.3%)', True),
    ('● 95% CI 폭 평균: 39.0%p  → "임계값 4°C 불확실성이 TARR을 39%p 변동시킴"', True),
    ('● CI 폭이 넓은 이유: 53~57°C S-curve 경사 구간이 N(55, 4²) 샘플링 범위 안에 포함됨', False),
    ('   → CI 폭은 "불안정"의 징표가 아니라, S-curve 민감 구간의 존재를 정량적으로 보여주는 결과임', False),
], 0.45, 6.10, 12.40, 1.20, size=11)

set_notes(s,
"""Monte Carlo 결과입니다.
왼쪽 히스토그램은 집계구별 TARR 중앙값의 분포를 보여줍니다.
평균 70.2%로, 대부분의 집계구가 60~90% 구간에 몰려 있습니다.

오른쪽은 95% CI 폭의 분포입니다.
평균 39%p로, 임계값을 정확히 모를 때 TARR 추정치가 약 40%p 정도 달라질 수 있음을 의미합니다.

여기서 CI 폭이 넓다고 해서 분석이 잘못된 것이 아닙니다.
앞서 본 S-curve에서 53~57°C 구간에 2°C만 차이 나도 TARR이 50%p 가까이 변하는데,
Monte Carlo 샘플이 바로 이 민감 구간을 통과하기 때문에 CI가 넓어진 것입니다.
이 자체가 "55°C 주변에서 임계값 선택이 결정적으로 중요하다"는 발견입니다.""")


# ═══════════════════════════════════════════════════════════════════════════
# 슬라이드 12: Chapter 4 — MC CI 공간지도
# ═══════════════════════════════════════════════════════════════════════════
s = new_slide()
add_chapter_header(s, 4, '분석 결과', 'Monte Carlo — 집계구별 공간 분포')

img_path = os.path.join(FIG_DIR, 'mc_tarr_ci_map.png')
if os.path.exists(img_path):
    add_image_fit(s, img_path, 0.40, 1.50, 12.50, 5.50,
                  caption='그림 6. (좌) 집계구별 TARR 중앙값 공간 분포 / (우) 95% CI 폭 공간 분포')

add_multiline(s, [
    ('● 좌: TARR 중앙값 — 한강변·서울숲 주변 고취약, 건물 밀집 구간 상대적 저취약', False),
    ('● 우: CI 폭 — S-curve 민감 구간에 걸리는 집계구일수록 CI 폭이 넓게 나타남', False),
], 0.45, 7.10, 12.40, 0.30, size=11)

set_notes(s,
"""이 그림은 Monte Carlo 결과를 공간적으로 표현한 것입니다.
왼쪽은 TARR 중앙값, 오른쪽은 95% CI 폭입니다.

CI 폭이 넓은 집계구를 보면,
임계값 변화에 특히 민감한 위치, 즉 S-curve 경사 구간에 위치한 링크 비율이 높은 곳입니다.
반대로 CI 폭이 좁은 집계구는 어떤 임계값을 쓰더라도 결과가 비슷하게 나오는 곳입니다.
이런 공간 패턴 자체도 정책적으로 의미가 있습니다.
CI 폭이 넓은 집계구일수록 임계값 결정이 접근성 분석 결과에 크게 영향을 미치므로,
더 정밀한 MRT 관측이 필요한 우선 지역으로 볼 수 있습니다.""")


# ═══════════════════════════════════════════════════════════════════════════
# 슬라이드 13: Chapter 4 — 3시간대 비교
# ═══════════════════════════════════════════════════════════════════════════
s = new_slide()
add_chapter_header(s, 4, '분석 결과', '시간대별 비교: 9시 · 13시 · 18시')

img_path = os.path.join(FIG_DIR, 'tarr_3hour_comparison.png')
if os.path.exists(img_path):
    add_image_fit(s, img_path, 0.40, 1.50, 12.50, 4.80,
                  caption='그림 7. 시간대별 MRT 최댓값, 고온링크 비율(MRT≥55°C), 집계구 평균 TARR 비교')

add_multiline(s, [
    ('● 9시: MRT 낮음 → 고온링크 비율·TARR 모두 낮음 — 폭염 이전 시간대', False),
    ('● 13시: MRT 최대 → 고온링크 비율 최대 → TARR 최대 (73.3%) — 분석 기준 시간대', True),
    ('● 18시: MRT 감소 시작 → TARR 감소 — 일몰 전 회복 단계', False),
    ('→ 13시가 접근성 제약이 가장 심각한 시간대임을 확인', True),
], 0.45, 6.50, 12.40, 0.90, size=11)

set_notes(s,
"""이 그림은 9시, 13시, 18시 세 시간대를 비교한 결과입니다.
세 패널 모두 13시에 가장 극단적인 값을 보입니다.
MRT가 가장 높고, 고온링크 비율도 가장 높으며, TARR도 가장 큽니다.
이것은 태양 복사가 최대인 13시에 폭염으로 인한 보행 접근성 제약이 가장 심각하다는 것을 보여줍니다.
9시는 아직 MRT가 낮아서 거의 제약이 없고,
18시는 조금씩 회복되는 단계입니다.
이 비교는 우리가 13시를 분석 기준 시간대로 선택한 이유를 뒷받침합니다.""")


# ═══════════════════════════════════════════════════════════════════════════
# 슬라이드 14: Chapter 4 — 역별 정류장 손실
# ═══════════════════════════════════════════════════════════════════════════
s = new_slide()
add_chapter_header(s, 4, '분석 결과', '지하철역별 접근 가능 정류장 손실율')

img_path = os.path.join(FIG_DIR, 'station_stop_loss_table.png')
if os.path.exists(img_path):
    add_image_fit(s, img_path, 0.40, 1.50, 12.50, 5.40,
                  caption='그림 8. 지하철역 7개소별 Classic·Thermal Catchment 정류장 수 및 손실율 (MRT 55°C, 13시)')

add_multiline(s, [
    ('● 서울숲역 80.0%, 행당역 77.6% 등 역별 손실율 차이 → 역 주변 열환경 이질성 반영', False),
    ('● 손실율이 높은 역일수록 폭염 시 환승·접근 대안이 더 제한됨', False),
], 0.45, 7.05, 12.40, 0.35, size=11)

set_notes(s,
"""이 표는 성동구 지하철역 7개소 각각에서 Thermal Catchment를 적용했을 때
접근 가능한 정류장이 얼마나 줄어드는지를 보여줍니다.

서울숲역은 80.0%, 행당역은 77.6%로 손실율이 높습니다.
반면 다른 역들은 상대적으로 낮은데,
이는 역 주변 도로망 구성과 건물 차폐 패턴이 다르기 때문입니다.

이런 역별 차이는 정책적으로 의미가 있습니다.
손실율이 높은 역일수록 폭염 시 환승이나 접근이 더 어렵고,
쿨링 인프라나 그늘막 등의 우선 설치 대상이 될 수 있습니다.""")


# ═══════════════════════════════════════════════════════════════════════════
# 슬라이드 15: Chapter 4 — OLS 회귀분석
# ═══════════════════════════════════════════════════════════════════════════
s = new_slide()
add_chapter_header(s, 4, '분석 결과', '공간 환경 변수 × TARR OLS 회귀분석 (보조 분석)')

img_path = os.path.join(FIG_DIR, 'regression_scatter.png')
if os.path.exists(img_path):
    add_image_fit(s, img_path, 0.40, 1.50, 12.50, 4.50,
                  caption='그림 9. 공간 환경 변수 5개 × TARR OLS 회귀 산점도 (n=506)')

add_multiline(s, [
    ('● OLS 결과: R² = 0.174 (모형 설명력 17.4%)', True),
    ('● hot_link_ratio (r=0.332***): 고온링크 비율이 높은 집계구일수록 TARR 높음 — 가장 강한 관계', False),
    ('● SVF (r=0.295***): SVF가 높은(개방된) 집계구일수록 TARR 높음 — 그늘 부족 ↔ 고온 링크', False),
    ('● H/W ratio (r=−0.208***): 협곡비가 높을수록 TARR 낮음 — 건물 차폐로 그늘 확보', False),
    ('● 해석: 이 회귀분석은 "MRT 선택의 근거"가 아닌 보조적 탐색 분석입니다.', True),
    ('   MRT 활용의 타당성은 SOLWEIG 선행연구(Lindberg & Grimmond 2011 등)로 뒷받침합니다.', False),
], 0.45, 6.20, 12.40, 1.20, size=11)

set_notes(s,
"""마지막 결과로 공간 환경 변수와 TARR의 관계를 OLS 회귀로 탐색했습니다.

전체 설명력 R²는 0.174로 높지 않지만,
방향성은 물리적으로 타당합니다.
고온링크 비율이 높은 집계구, SVF가 높은 개방된 집계구에서 TARR이 높고,
건물 협곡비가 높아 그늘이 많은 곳에서는 TARR이 낮습니다.

이 회귀분석은 MRT 방법론의 근거가 아니라 보조적 탐색입니다.
MRT 활용의 타당성은 SOLWEIG 등 기존 선행연구로 뒷받침하고,
이 분석은 "어떤 공간 환경이 접근성 손실을 더 크게 만드는가"를 보여주는 부수 결과입니다.""")


# ═══════════════════════════════════════════════════════════════════════════
# 슬라이드 16: Chapter 5 — 한계점 및 향후 연구
# ═══════════════════════════════════════════════════════════════════════════
s = new_slide()
add_chapter_header(s, 5, '한계점 및 향후 연구', '한계점')

add_multiline(s, [
    ('● 바람 속도 균일 가정: va = 2.37 m/s 전 구간 동일 적용', False),
    ('   → 도시 협곡 내 실제 풍속은 2~5배 변이 → 대류 효과 과소·과대 혼재', False),
    ('   → 향후: CFD 시뮬레이션 또는 다점 S-DoT 풍속 보정 필요', False),
    ('', False),
    ('● GHI 단일 지점: 성동구 전체에 동일 일사량 적용', False),
    ('   → 구름 패턴 등 미세 공간 차이 미반영', False),
    ('', False),
    ('● MRT 실측 검증 부재: SOLWEIG 약식 추정값 — 복사계 실측 비교 미수행', False),
    ('', False),
    ('● Hard Cut 행동 가정: 보행자가 고온 링크를 완전 회피한다는 실증 근거 없음', False),
    ('   → "물리적 열스트레스 임계값"으로 재해석하여 행동적 의미와 구분', False),
    ('', False),
    ('● 집계구 도심 가정: 실제 거주민은 집계구 전역 분포 → 도심 출발은 근사치', False),
    ('', False),
    ('● 분석 기간: 2025.07.28–08.03 단일 폭염 이벤트(7일) — 사례 연구 수준, 일반화 한계', False),
    ('', False),
    ('● OSM 네트워크: 지하도·실내 공조 구간 등 열환경 양호 경로 미포함', False),
], 0.45, 1.88, 6.20, 5.20, size=11)

add_section_box(s, '향후 연구 방향', 1.88)

add_multiline(s, [
    ('향후 연구 방향', True),
    ('', False),
    ('● 고령자 보행 속도(3.0 km/h) 분리 적용', False),
    ('   → 폭염 취약계층(노인) 접근성 분석', False),
    ('', False),
    ('● QWEN 멀티모달 LLM 거리뷰 분석', False),
    ('   → 거리 영상에서 SVF·그림자 자동 추출', False),
    ('   → MRT 추정 자동화 파이프라인', False),
    ('', False),
    ('● CFD 또는 다점 S-DoT 풍속 연계', False),
    ('   → 바람 공간 변이 반영한 MRT 정밀화', False),
    ('', False),
    ('● 쿨링 인프라 최적 배치 분석', False),
    ('   → TARR 고취약 집계구 대상', False),
    ('   → 그늘막·수경시설 배치 시뮬레이션', False),
], 6.60, 1.88, 6.50, 5.20, size=11)

set_notes(s,
"""연구의 한계점을 솔직하게 밝히겠습니다.

가장 큰 한계는 바람 속도를 전 구간 동일하게 적용한 것입니다.
실제 도시에서는 건물 배치에 따라 풍속이 크게 달라지는데, 이를 반영하지 못했습니다.

또한 MRT를 실제 측정값과 비교 검증하지 못한 점,
그리고 분석 기간이 7일 단일 폭염 이벤트에 한정된 점도 한계입니다.

향후 연구로는 고령자 보행 속도를 분리해서 취약계층 분석,
QWEN 같은 멀티모달 AI를 이용한 거리뷰 영상 기반 MRT 추정 자동화,
그리고 TARR이 높은 집계구에 쿨링 인프라를 최적 배치하는 시뮬레이션 등을 계획하고 있습니다.""")


# ═══════════════════════════════════════════════════════════════════════════
# 슬라이드 17: Chapter 6 — 결론
# ═══════════════════════════════════════════════════════════════════════════
s = new_slide()
add_chapter_header(s, 6, '결론', '요약 및 기여')

add_multiline(s, [
    ('● 개념적 기여: Thermal Catchment Area 제안', True),
    ('   → 폭염 시 MRT 기반 Hard Cut을 적용한 새로운 보행 접근성 공간 단위', False),
    ('   → 기존 Catchment가 무시해온 열환경 제약을 공간 단위 자체에 내재화', False),
    ('', False),
    ('● 실증 결과: 성동구 집계구(n=506), MRT 55°C 기준, 13시 기준', True),
    ('   → 집계구 평균 TARR 73.3% — 폭염 시 약 73%의 정류장 접근 기회 손실', False),
    ('   → 서울숲역 80.0%, 행당역 77.6% 등 역별 손실율 차이 확인', False),
    ('', False),
    ('● 방법론적 기여: S-curve 발견 + Monte Carlo 불확실성 정량화', True),
    ('   → 53~57°C 구간에서 TARR 88%→39% 급변 — 임계값 선택이 결과에 결정적', False),
    ('   → Monte Carlo (N=2,000)로 임계값 불확실성을 확률적으로 전파', False),
    ('   → 단일값 결과(73.3%)가 Monte Carlo 중앙값(70.2%)과 근접 — robust 확인', False),
    ('', False),
    ('● 정책적 함의', True),
    ('   → TARR 고취약 집계구를 우선 쿨링 인프라 정비 대상으로 공간적 식별 가능', False),
    ('   → Thermal Catchment는 기후 적응형 도시 접근성 계획의 공간 기준 단위로 활용 가능', False),
    ('', False),
], 0.45, 1.88, 12.40, 5.00, size=12)

add_rect(s, 0.40, 7.10, 12.53, 0.02, C_RED)

set_notes(s,
"""결론을 말씀드리겠습니다.

이 연구는 세 가지 기여를 합니다.

첫째, Thermal Catchment Area라는 새로운 공간 단위 개념을 제안했습니다.
폭염 때 MRT가 높은 길을 피하면 실제로 갈 수 있는 정류장이 줄어드는데,
이것을 공간 단위로 정의하고 측정한 것입니다.

둘째, 서울 성동구에서 실증한 결과, 폭염 피크 시간인 13시에 평균 73%의 정류장 접근 기회가 사라졌습니다.
역별로도 차이가 있어서 서울숲역이 가장 심각합니다.

셋째, S-curve라는 중요한 발견과 Monte Carlo로 불확실성을 정량화했습니다.
53~57°C 구간이 임계값 선택의 변곡점이며,
Monte Carlo 결과가 단일값 결과와 근접하여 55°C 기준의 robustness를 확인했습니다.

정책적으로는 TARR이 높은 집계구를 우선 쿨링 인프라 설치 대상으로 식별하고,
기후 적응형 도시 계획의 공간 기준으로 활용할 수 있습니다.

감사합니다.""")


# ── 저장 ──────────────────────────────────────────────────────────────────
prs.save(OUT_PATH)
print(f"저장: {OUT_PATH}")
print(f"슬라이드 수: {len(prs.slides)}")
